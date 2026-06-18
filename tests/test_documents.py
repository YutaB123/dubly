"""Tests for the make-a-document (Word .docx) tool."""

from __future__ import annotations

import io
import zipfile

from app.documents import (
    DocumentService, render_docx, render_pdf, render_pptx, render_xlsx, render_csv,
)
from app.db import FileStore

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


class FakeSms:
    def __init__(self):
        self.sent = []

    def send(self, text, to=None, media_url=None):
        self.sent.append((text, media_url))


class FakeOneDrive:
    def __init__(self):
        self.uploaded = []

    def upload(self, name, data, content_type):
        self.uploaded.append((name, content_type, len(data)))


def test_render_pdf_returns_pdf_bytes_and_handles_unicode():
    data = render_pdf("My Title", "first line\n\nsecond — with “smart” quotes and …")
    assert data[:4] == b"%PDF"
    assert len(data) > 200


def test_render_docx_is_a_valid_word_zip_with_the_text():
    data = render_docx("My Title", "Hello world\n\nSecond paragraph")
    assert data[:2] == b"PK"  # zip magic
    z = zipfile.ZipFile(io.BytesIO(data))
    names = z.namelist()
    assert "[Content_Types].xml" in names
    assert "word/document.xml" in names
    doc = z.read("word/document.xml").decode("utf-8")
    assert "My Title" in doc
    assert "Hello world" in doc
    assert "Second paragraph" in doc


def test_render_docx_escapes_xml_special_characters():
    data = render_docx("T", "a < b & c > d")
    doc = zipfile.ZipFile(io.BytesIO(data)).read("word/document.xml").decode("utf-8")
    assert "&lt;" in doc and "&amp;" in doc and "&gt;" in doc


def test_make_document_sends_downloadable_docx_and_copies_to_onedrive(tmp_path):
    sms = FakeSms()
    files = FileStore(tmp_path / "f.sqlite")
    od = FakeOneDrive()
    svc = DocumentService(sms=sms, files=files, public_base_url="https://app.example", onedrive=od)

    out = svc.dispatch("make_document", {
        "title": "STAT 311 Study Guide",
        "content": "point one\npoint two",
    })
    assert "sent" in out.lower()

    # Sent as a media link.
    text, media_url = sms.sent[0]
    assert "STAT 311 Study Guide" in text
    assert media_url[0].startswith("https://app.example/file/")

    # Stored as a servable Word document.
    fid = media_url[0].rsplit("/", 1)[1]
    filename, ctype, data = files.get(fid)
    assert filename == "STAT 311 Study Guide.docx"
    assert ctype == DOCX_MIME
    assert data[:2] == b"PK"

    # Also copied into the OneDrive folder.
    assert od.uploaded and od.uploaded[0][0] == "STAT 311 Study Guide.docx"


def test_render_pptx_makes_a_valid_deck_with_the_content():
    data = render_pptx("STAT 311 Guide", "Key Ideas:\nfirst point\nsecond point")
    assert data[:2] == b"PK"
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    blob = " ".join(sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame)
    assert "STAT 311 Guide" in blob
    assert "first point" in blob and "second point" in blob


def _doc_svc(tmp_path):
    sms = FakeSms()
    files = FileStore(tmp_path / "f.sqlite")
    svc = DocumentService(sms=sms, files=files, public_base_url="https://app.example", onedrive=None)
    return svc, sms, files


def _stored(sms, files):
    fid = sms.sent[0][1][0].rsplit("/", 1)[1]
    return files.get(fid)


def test_make_document_pdf_format(tmp_path):
    svc, sms, files = _doc_svc(tmp_path)
    svc.dispatch("make_document", {"title": "Notes", "content": "hi", "format": "pdf"})
    filename, ctype, data = _stored(sms, files)
    assert filename == "Notes.pdf"
    assert ctype == "application/pdf"
    assert data[:4] == b"%PDF"


def test_make_document_powerpoint_format(tmp_path):
    svc, sms, files = _doc_svc(tmp_path)
    svc.dispatch("make_document", {"title": "Deck", "content": "Topic:\npoint", "format": "powerpoint"})
    filename, ctype, data = _stored(sms, files)
    assert filename == "Deck.pptx"
    assert ctype == PPTX_MIME
    assert data[:2] == b"PK"


def test_render_csv_is_the_raw_content_bytes():
    assert render_csv("Grades", "name,grade\nDubs,A") == b"name,grade\nDubs,A"


def test_render_xlsx_is_a_valid_workbook_with_the_data():
    import io, zipfile
    data = render_xlsx("Grades", "name,grade\nDubs,95")
    assert data[:2] == b"PK"
    z = zipfile.ZipFile(io.BytesIO(data))
    assert "xl/workbook.xml" in z.namelist()
    blob = b"".join(z.read(n) for n in z.namelist() if n.endswith(".xml"))
    assert b"name" in blob and b"Dubs" in blob


def test_make_document_excel_format(tmp_path):
    svc, sms, files = _doc_svc(tmp_path)
    svc.dispatch("make_document", {"title": "Data", "content": "a,b\n1,2", "format": "excel"})
    filename, ctype, data = _stored(sms, files)
    assert filename == "Data.xlsx" and ctype == XLSX_MIME and data[:2] == b"PK"


def test_make_document_csv_format(tmp_path):
    svc, sms, files = _doc_svc(tmp_path)
    svc.dispatch("make_document", {"title": "Rows", "content": "a,b\n1,2", "format": "csv"})
    filename, ctype, data = _stored(sms, files)
    assert filename == "Rows.csv" and ctype == "text/csv" and data == b"a,b\n1,2"


def test_make_document_defaults_to_word(tmp_path):
    svc, sms, files = _doc_svc(tmp_path)
    svc.dispatch("make_document", {"title": "Plain", "content": "x"})
    filename, ctype, _ = _stored(sms, files)
    assert filename == "Plain.docx" and ctype == DOCX_MIME


def test_make_document_works_without_onedrive(tmp_path):
    sms = FakeSms()
    svc = DocumentService(sms=sms, files=FileStore(tmp_path / "f.sqlite"),
                          public_base_url="https://app.example", onedrive=None)
    svc.dispatch("make_document", {"title": "Notes", "content": "stuff"})
    assert sms.sent and sms.sent[0][1][0].startswith("https://app.example/file/")
