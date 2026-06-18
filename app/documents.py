"""Make a downloadable Word (.docx) document from text and send it to the student.

The brain writes the content; this turns it into a real .docx (openable/editable in
Word or Google Docs), sends it as a downloadable link, and drops a copy in OneDrive.
"""

from __future__ import annotations

import io
import re
import uuid
import zipfile
from xml.sax.saxutils import escape

from fpdf import FPDF
from fpdf.enums import XPos, YPos

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

DOCUMENT_TOOLS = [
    {
        "name": "make_document",
        "description": "Create a downloadable file from text and SEND it to the student. Use "
        "whenever they ask you to write something up as a file they can download — a study guide, "
        "summary, notes, outline, cheat sheet, essay blueprint, slides, table of data, etc. You "
        "write the full content yourself in 'content'. Pick 'format' to match what they ask for: "
        "'word' (.docx, default), 'pdf', 'powerpoint' (.pptx slides), 'excel' (.xlsx spreadsheet), "
        "or 'csv'. For powerpoint, write content as blank-line-separated sections, each a short "
        "heading line then bullet lines (one slide per section). For excel/csv, write content as "
        "comma-separated rows, one per line, first row as column headers.",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "Short title / filename, no extension (e.g. 'STAT 311 Study Guide').",
                },
                "content": {
                    "type": "string",
                    "description": "The full document text. Plain text; blank lines between paragraphs/sections.",
                },
                "format": {
                    "type": "string",
                    "enum": ["word", "pdf", "powerpoint", "excel", "csv"],
                    "description": "File type to produce. Default 'word'.",
                },
            },
            "required": ["title", "content"],
        },
    }
]

# fpdf2 core fonts are latin-1; map common unicode so content doesn't get mangled.
_UNICODE = {
    "—": "-", "–": "-", "’": "'", "‘": "'",
    "“": '"', "”": '"', "…": "...", "•": "-",
    " ": " ", "→": "->", "←": "<-",
}


def _latin1(text: str) -> str:
    for k, v in _UNICODE.items():
        text = text.replace(k, v)
    return text.encode("latin-1", "replace").decode("latin-1")


def render_pdf(title: str, content: str) -> bytes:
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(True, margin=15)
    pdf.set_font("Helvetica", "B", 16)
    pdf.multi_cell(0, 10, _latin1(title), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.ln(2)
    pdf.set_font("Helvetica", "", 12)
    for line in content.split("\n"):
        pdf.multi_cell(0, 7, _latin1(line) or " ", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    return bytes(pdf.output())


# --- Word .docx (a minimal, valid OOXML package, stdlib only) ----------------

_CONTENT_TYPES = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
    '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
    '<Default Extension="xml" ContentType="application/xml"/>'
    '<Override PartName="/word/document.xml" '
    'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
    "</Types>"
)

_RELS = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
    '<Relationship Id="rId1" '
    'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
    'Target="word/document.xml"/></Relationships>'
)

_W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"


def _paragraph(text: str, bold: bool = False, size: int | None = None) -> str:
    rpr = ""
    if bold or size:
        rpr = "<w:rPr>" + ("<w:b/>" if bold else "") + (
            f'<w:sz w:val="{size}"/>' if size else ""
        ) + "</w:rPr>"
    return f'<w:p><w:r>{rpr}<w:t xml:space="preserve">{escape(text)}</w:t></w:r></w:p>'


def render_docx(title: str, content: str) -> bytes:
    """A real Word .docx: bold title, then one paragraph per line of content."""
    paras = [_paragraph(title, bold=True, size=32)]
    paras += [_paragraph(line) for line in content.split("\n")]
    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<w:document xmlns:w="{_W_NS}"><w:body>{"".join(paras)}</w:body></w:document>'
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", _CONTENT_TYPES)
        z.writestr("_rels/.rels", _RELS)
        z.writestr("word/document.xml", document_xml)
    return buf.getvalue()


# --- PowerPoint (.pptx) via python-pptx --------------------------------------

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def render_pptx(title: str, content: str) -> bytes:
    """A slide deck: a title slide, then one slide per blank-line-separated section
    (first line of each section is the slide heading, the rest become bullets)."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "made by Dubly 🐾"

    sections = [s for s in content.split("\n\n") if s.strip()]
    if not sections:
        sections = [content] if content.strip() else []
    for sec in sections:
        lines = [ln.strip() for ln in sec.splitlines() if ln.strip()]
        if not lines:
            continue
        heading = lines[0].rstrip(":")
        bullets = [ln.lstrip("-•* ") for ln in lines[1:]] or [heading]
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = heading
        body = slide.placeholders[1].text_frame
        body.text = bullets[0]
        for b in bullets[1:]:
            p = body.add_paragraph()
            p.text = b
        for p in body.paragraphs:
            for run in p.runs:
                run.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- CSV + Excel (.xlsx) -----------------------------------------------------

XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


def render_csv(title: str, content: str) -> bytes:
    """A CSV file — the content is already comma-separated rows, just encode it."""
    return (content or "").encode("utf-8")


def render_xlsx(title: str, content: str) -> bytes:
    """An Excel workbook from comma-separated rows; row 1 is treated as bold headers,
    and numeric cells are written as numbers."""
    import csv as _csv
    import xlsxwriter

    buf = io.BytesIO()
    wb = xlsxwriter.Workbook(buf, {"in_memory": True})
    sheet_name = re.sub(r"[\[\]:*?/\\]", "", title)[:31] or "Sheet1"
    ws = wb.add_worksheet(sheet_name)
    bold = wb.add_format({"bold": True})
    rows = list(_csv.reader(io.StringIO(content or "")))
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            if r == 0:
                ws.write(r, c, val, bold)
            else:
                try:
                    ws.write_number(r, c, float(val))
                except (ValueError, TypeError):
                    ws.write(r, c, val)
    wb.close()
    return buf.getvalue()


# extension + MIME for each format the document tool can produce.
_FORMATS = {
    "pdf": (".pdf", "application/pdf"),
    "powerpoint": (".pptx", PPTX_MIME),
    "pptx": (".pptx", PPTX_MIME),
    "slides": (".pptx", PPTX_MIME),
    "presentation": (".pptx", PPTX_MIME),
    "excel": (".xlsx", XLSX_MIME),
    "xlsx": (".xlsx", XLSX_MIME),
    "spreadsheet": (".xlsx", XLSX_MIME),
    "sheet": (".xlsx", XLSX_MIME),
    "csv": (".csv", "text/csv"),
    "word": (".docx", DOCX_MIME),
    "docx": (".docx", DOCX_MIME),
}


def _safe_name(title: str) -> str:
    return re.sub(r"[^\w\- ]", "", title).strip() or "document"


class DocumentService:
    def __init__(self, sms, files, public_base_url: str, onedrive=None):
        self.sms = sms
        self.files = files
        self.public_base_url = public_base_url.rstrip("/")
        self.onedrive = onedrive

    def tool_names(self) -> list[str]:
        return [t["name"] for t in DOCUMENT_TOOLS]

    def schemas(self) -> list[dict]:
        return list(DOCUMENT_TOOLS)

    def dispatch(self, name: str, tool_input: dict) -> str:
        if name != "make_document":
            return f"(unknown document tool: {name})"
        title = (tool_input.get("title") or "document").strip()
        content = tool_input.get("content") or ""
        fmt = (tool_input.get("format") or "word").strip().lower()
        ext, mime = _FORMATS.get(fmt, _FORMATS["word"])
        if ext == ".pdf":
            data = render_pdf(title, content)
        elif ext == ".pptx":
            data = render_pptx(title, content)
        elif ext == ".xlsx":
            data = render_xlsx(title, content)
        elif ext == ".csv":
            data = render_csv(title, content)
        else:
            data = render_docx(title, content)
        filename = _safe_name(title) + ext
        file_id = uuid.uuid4().hex
        self.files.save(file_id, filename, mime, data)
        # Also drop a copy in their OneDrive folder (so it's on the laptop too).
        if self.onedrive is not None:
            try:
                self.onedrive.upload(filename, data, mime)
            except Exception:
                pass
        self.sms.send(
            f"here's {title}:", media_url=[f"{self.public_base_url}/file/{file_id}"]
        )
        return f"created and sent {filename} ({len(data)} bytes) to them."
